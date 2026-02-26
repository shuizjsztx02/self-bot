#!/usr/bin/env python3
"""
创建关于 Qwen 大模型的 PPT 演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_qwen_presentation(output_path="test_output/qwen_intro.pptx"):
    """创建 Qwen 大模型 PPT"""
    
    # 创建演示文稿
    prs = Presentation()
    
    # 设置幻灯片尺寸为 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # ====== 第1页：封面 ======
    slide_layout = prs.slide_layouts[0]  # 标题幻灯片
    slide1 = prs.slides.add_slide(slide_layout)
    
    # 设置背景色
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(30, 39, 97)  # 午夜蓝
    
    # 标题
    title = slide1.shapes.title
    title.text = "Qwen 大模型介绍"
    title.text_frame.paragraphs[0].font.size = Pt(48)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle = slide1.placeholders[1]
    subtitle.text = "通义千问 - 阿里巴巴集团研发的先进大语言模型"
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(202, 220, 252)  # 冰蓝色
    subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ====== 第2页：Qwen 模型概述 ======
    slide_layout = prs.slide_layouts[1]  # 标题和内容
    slide2 = prs.slides.add_slide(slide_layout)
    
    # 标题
    title = slide2.shapes.title
    title.text = "Qwen 模型概述"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 39, 97)
    
    # 内容
    content = slide2.placeholders[1]
    text_frame = content.text_frame
    text_frame.clear()  # 清除默认文本
    
    # 添加要点
    p = text_frame.add_paragraph()
    p.text = "• 开发团队：阿里巴巴集团达摩院"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 0, 0)
    
    p = text_frame.add_paragraph()
    p.text = "• 模型系列：Qwen-1.5、Qwen-2、Qwen-2.5、Qwen-2.5-Coder 等"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 0, 0)
    
    p = text_frame.add_paragraph()
    p.text = "• 参数规模：0.5B 到 72B 多种规格"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 0, 0)
    
    p = text_frame.add_paragraph()
    p.text = "• 开源协议：Apache 2.0 开源"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 0, 0)
    
    p = text_frame.add_paragraph()
    p.text = "• 主要特点：多语言支持、代码生成、数学推理"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 0, 0)
    
    # ====== 第3页：主要特性 ======
    slide_layout = prs.slide_layouts[1]
    slide3 = prs.slides.add_slide(slide_layout)
    
    # 标题
    title = slide3.shapes.title
    title.text = "主要特性"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 39, 97)
    
    # 内容 - 使用两列布局
    content = slide3.placeholders[1]
    text_frame = content.text_frame
    text_frame.clear()
    
    # 左列特性
    p = text_frame.add_paragraph()
    p.text = "🔹 强大的多语言能力"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(30, 39, 97)
    p.font.bold = True
    
    p = text_frame.add_paragraph()
    p.text = "  支持中、英、日、韩等多种语言"
    p.font.size = Pt(18)
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "🔹 卓越的代码生成"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(30, 39, 97)
    p.font.bold = True
    
    p = text_frame.add_paragraph()
    p.text = "  支持 Python、Java、C++ 等多种编程语言"
    p.font.size = Pt(18)
    p.level = 1
    
    # 右列特性
    p = text_frame.add_paragraph()
    p.text = "🔹 数学推理能力"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(30, 39, 97)
    p.font.bold = True
    
    p = text_frame.add_paragraph()
    p.text = "  在数学竞赛和推理任务中表现优异"
    p.font.size = Pt(18)
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "🔹 长上下文支持"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(30, 39, 97)
    p.font.bold = True
    
    p = text_frame.add_paragraph()
    p.text = "  支持 32K 甚至 128K 长上下文"
    p.font.size = Pt(18)
    p.level = 1
    
    # ====== 第4页：应用场景 ======
    slide_layout = prs.slide_layouts[1]
    slide4 = prs.slides.add_slide(slide_layout)
    
    # 标题
    title = slide4.shapes.title
    title.text = "应用场景"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 39, 97)
    
    # 内容
    content = slide4.placeholders[1]
    text_frame = content.text_frame
    text_frame.clear()
    
    # 添加应用场景
    scenarios = [
        ("💼 企业应用", "智能客服、文档分析、代码助手"),
        ("🎓 教育领域", "个性化辅导、作业批改、学习伙伴"),
        ("🔬 科研工作", "文献总结、实验设计、数据分析"),
        ("💻 开发工具", "代码生成、调试助手、API文档"),
        ("📱 个人助手", "内容创作、翻译、知识问答")
    ]
    
    for scenario, description in scenarios:
        p = text_frame.add_paragraph()
        p.text = f"• {scenario}"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(30, 39, 97)
        p.font.bold = True
        
        p = text_frame.add_paragraph()
        p.text = f"  {description}"
        p.font.size = Pt(18)
        p.level = 1
    
    # ====== 第5页：总结 ======
    slide_layout = prs.slide_layouts[0]  # 标题幻灯片
    slide5 = prs.slides.add_slide(slide_layout)
    
    # 设置背景色
    background = slide5.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(30, 39, 97)  # 午夜蓝
    
    # 标题
    title = slide5.shapes.title
    title.text = "总结"
    title.text_frame.paragraphs[0].font.size = Pt(48)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 内容
    content = slide5.placeholders[1]
    text_frame = content.text_frame
    text_frame.clear()
    
    summary_points = [
        "Qwen 是阿里巴巴研发的先进大语言模型",
        "开源、多语言、多模态支持",
        "在代码生成和数学推理方面表现突出",
        "适用于多种商业和学术场景",
        "持续更新，社区活跃"
    ]
    
    for point in summary_points:
        p = text_frame.add_paragraph()
        p.text = f"✓ {point}"
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(202, 220, 252)  # 冰蓝色
        p.alignment = PP_ALIGN.CENTER
    
    # 保存演示文稿
    prs.save(output_path)
    print(f"PPT 已保存到: {output_path}")
    print(f"总页数: {len(prs.slides)}")
    
    return output_path

if __name__ == "__main__":
    output_path = create_qwen_presentation()
    print(f"演示文稿创建完成: {output_path}")