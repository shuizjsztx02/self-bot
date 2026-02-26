from typing import List, Dict, Any, Optional
from pathlib import Path
import asyncio

from .base import DocumentParser, ParsedDocument, ChunkResult


class PPTXParser(DocumentParser):
    
    async def parse(self, file_path: str) -> ParsedDocument:
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, self._parse_sync, file_path)
    
    def _parse_sync(self, file_path: str) -> ParsedDocument:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx not installed. Run: pip install python-pptx")
        
        prs = Presentation(file_path)
        
        slides = []
        all_text = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            title = None
            
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text = shape.text.strip()
                    slide_text.append(text)
                    
                    if shape.has_text_frame:
                        if shape.text_frame.paragraphs:
                            first_para = shape.text_frame.paragraphs[0]
                            if first_para.text.strip():
                                title = first_para.text.strip()
            
            slide_content = '\n'.join(slide_text)
            all_text.append(slide_content)
            
            slides.append({
                'slide_number': slide_num,
                'title': title,
                'content': slide_content,
            })
        
        content = '\n\n--- 幻灯片 ---\n\n'.join(all_text)
        
        metadata = {
            'filename': Path(file_path).name,
            'file_type': 'pptx',
            'slide_count': len(prs.slides),
        }
        
        return ParsedDocument(
            content=content,
            metadata=metadata,
            pages=slides,
        )
    
    def supported_extensions(self) -> List[str]:
        return ['.pptx']
    
    def chunk_by_slides(
        self,
        parsed_doc: ParsedDocument,
        slides_per_chunk: int = 3,
        chunk_size: int = 1000,
    ) -> List[ChunkResult]:
        chunks = []
        
        if not parsed_doc.pages:
            text_chunks = self.chunk_text(parsed_doc.content, chunk_size)
            for i, chunk in enumerate(text_chunks):
                chunks.append(ChunkResult(
                    content=chunk,
                    token_count=self.count_tokens(chunk),
                    metadata={'chunk_index': i},
                ))
            return chunks
        
        current_chunk_slides = []
        current_chunk_text = ""
        
        for slide in parsed_doc.pages:
            slide_text = slide['content']
            
            if len(current_chunk_text) + len(slide_text) > chunk_size and current_chunk_text:
                chunks.append(ChunkResult(
                    content=current_chunk_text.strip(),
                    token_count=self.count_tokens(current_chunk_text),
                    page_number=current_chunk_slides[0] if current_chunk_slides else None,
                    metadata={
                        'slide_range': f"{current_chunk_slides[0]}-{current_chunk_slides[-1]}" if len(current_chunk_slides) > 1 else str(current_chunk_slides[0]),
                    },
                ))
                current_chunk_text = ""
                current_chunk_slides = []
            
            current_chunk_text += ('\n\n' if current_chunk_text else '') + slide_text
            current_chunk_slides.append(slide['slide_number'])
        
        if current_chunk_text:
            chunks.append(ChunkResult(
                content=current_chunk_text.strip(),
                token_count=self.count_tokens(current_chunk_text),
                page_number=current_chunk_slides[0] if current_chunk_slides else None,
                metadata={
                    'slide_range': f"{current_chunk_slides[0]}-{current_chunk_slides[-1]}" if len(current_chunk_slides) > 1 else str(current_chunk_slides[0]),
                },
            ))
        
        return chunks
